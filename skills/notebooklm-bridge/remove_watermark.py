#!/usr/bin/env python3
"""
NotebookLM Watermark Remover

Removes the NotebookLM logo watermark from PDF slides and infographics.
Two modes:
  (default)   Stamp page numbers over the watermark region (fast, clean)
  --inpaint   Use biharmonic inpainting to erase the watermark region

Usage:
    python remove_watermark.py input.pdf                    # page numbers (default)
    python remove_watermark.py input.pdf --stamp "My Title" # custom text on every page
    python remove_watermark.py input.pdf --inpaint          # biharmonic inpaint
    python remove_watermark.py input.pdf -o output.pdf      # custom output path
    python remove_watermark.py input.pdf --pptx             # PowerPoint file
    python remove_watermark.py input.pdf --png              # ZIP of clean PNGs
    python remove_watermark.py *.pdf                        # batch mode

Based on: https://huggingface.co/spaces/dseditor/WaterMarkLM (MIT)
"""

import argparse
import datetime
import json
import re
import sys
import tempfile
import time
import zipfile
from pathlib import Path
from urllib.request import urlopen

import fitz  # PyMuPDF
import numpy as np
from PIL import Image, ImageDraw, ImageFont
from skimage.restoration import inpaint


# Fallback coordinates (known-good as of 2026-04)
_FALLBACK_BASE_WIDTH = 2867
_FALLBACK_BASE_HEIGHT = 1600
_FALLBACK_LOGO_COORDS = (1530, 1595, 2620, 2860)  # (row1, row2, col1, col2)

# Auto-update from HuggingFace
_HF_APP_URL = "https://huggingface.co/spaces/dseditor/WaterMarkLM/raw/main/app.py"
_CACHE_FILE = Path(__file__).parent / ".watermark_coords_cache.json"
_CACHE_TTL = 86400 * 7  # refresh weekly

DPI = 150

# Font search paths (cross-platform)
_FONT_PATHS = [
    "C:/Windows/Fonts/segoeuib.ttf",    # Windows: Segoe UI Bold
    "C:/Windows/Fonts/segoeui.ttf",
    "C:/Windows/Fonts/arialbd.ttf",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",  # Linux
    "/System/Library/Fonts/Helvetica.ttc",  # macOS
]


def _fetch_coords_from_hf() -> tuple[int, int, int, tuple] | None:
    """Fetch latest watermark coordinates from the HuggingFace source."""
    try:
        resp = urlopen(_HF_APP_URL, timeout=10)
        code = resp.read().decode("utf-8")

        bw = re.search(r"self\.base_width\s*=\s*(\d+)", code)
        bh = re.search(r"self\.base_height\s*=\s*(\d+)", code)
        lc = re.search(r"self\.logo_coords\s*=\s*\(([^)]+)\)", code)

        if bw and bh and lc:
            coords = tuple(int(x.strip()) for x in lc.group(1).split(","))
            return int(bw.group(1)), int(bh.group(1)), coords
    except Exception:
        pass
    return None


def _load_coords() -> tuple[int, int, tuple]:
    """Load coords from cache or fetch from HuggingFace. Falls back to hardcoded."""
    # Try cache first
    if _CACHE_FILE.exists():
        try:
            data = json.loads(_CACHE_FILE.read_text(encoding="utf-8"))
            if time.time() - data.get("ts", 0) < _CACHE_TTL:
                return data["bw"], data["bh"], tuple(data["coords"])
        except Exception:
            pass

    # Fetch fresh
    result = _fetch_coords_from_hf()
    if result:
        bw, bh, coords = result
        try:
            _CACHE_FILE.write_text(
                json.dumps({"bw": bw, "bh": bh, "coords": list(coords), "ts": time.time()}),
                encoding="utf-8",
            )
        except Exception:
            pass
        return bw, bh, coords

    # Fallback
    return _FALLBACK_BASE_WIDTH, _FALLBACK_BASE_HEIGHT, _FALLBACK_LOGO_COORDS


BASE_WIDTH, BASE_HEIGHT, LOGO_COORDS = _load_coords()


def _load_font(size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    """Find and load a font, falling back to default."""
    import os
    for fp in _FONT_PATHS:
        if os.path.exists(fp):
            try:
                return ImageFont.truetype(fp, size)
            except Exception:
                continue
    return ImageFont.load_default()


def pdf_to_images(pdf_path: str, dpi: int = DPI) -> list[Path]:
    """Convert PDF pages to PNG images."""
    doc = fitz.open(pdf_path)
    tmp_dir = Path(tempfile.mkdtemp(prefix="watermark_"))
    images = []
    for i, page in enumerate(doc):
        mat = fitz.Matrix(dpi / 72, dpi / 72)
        pix = page.get_pixmap(matrix=mat)
        out = tmp_dir / f"page_{i+1:04d}.png"
        pix.save(str(out))
        images.append(out)
    doc.close()
    return images


def _get_watermark_region(image_path: Path) -> tuple[Image.Image, int, int, int, int] | None:
    """Get the watermark region coordinates scaled to the image size."""
    image = Image.open(image_path)
    w, h = image.size  # PIL uses (w, h)

    ratio = w / BASE_WIDTH
    r1, r2, c1, c2 = LOGO_COORDS
    r1, r2, c1, c2 = int(r1 * ratio), int(r2 * ratio), int(c1 * ratio), int(c2 * ratio)

    # Clamp to image bounds
    r1, r2 = max(0, r1), min(h, r2)
    c1, c2 = max(0, c1), min(w, c2)

    if r2 <= r1 or c2 <= c1:
        return None
    return image, r1, r2, c1, c2


def _detect_bg(image_array: np.ndarray, r1: int, r2: int, c1: int, c2: int):
    """Sample background color and brightness from strip above the watermark."""
    region_h = r2 - r1
    margin = max(5, int(region_h * 0.1))
    sample_r1 = max(0, r1 - margin)
    if sample_r1 < r1:
        bg_strip = image_array[sample_r1:r1, c1:c2]
        bg_color = np.mean(bg_strip, axis=(0, 1)).astype(np.uint8)
        brightness = (int(bg_color[0]) * 299 + int(bg_color[1]) * 587 + int(bg_color[2]) * 114) // 1000
    else:
        bg_color = np.array([30, 30, 30, 255], dtype=np.uint8)
        brightness = 30
    return bg_color, brightness


def detect_text_color(image_path: Path) -> tuple[int, int, int, int]:
    """Detect whether the watermark region background is light or dark.

    Returns the appropriate text color (dark for light bg, light for dark bg).
    """
    result = _get_watermark_region(image_path)
    if result is None:
        return (80, 80, 80, 200)  # default dark

    image, r1, r2, c1, c2 = result
    if image.mode != "RGBA":
        image = image.convert("RGBA")
    image_array = np.array(image)
    _, brightness = _detect_bg(image_array, r1, r2, c1, c2)

    if brightness > 127:
        return (80, 80, 80, 200)    # dark text on light bg
    else:
        return (200, 200, 200, 200)  # light text on dark bg


def stamp_text(
    image_path: Path,
    output_path: Path,
    text: str | None,
    text_color: tuple[int, int, int, int] | None = None,
) -> bool:
    """Replace the watermark region with stamped text (e.g. page number).

    If text is None, just fills the watermark region with background color.
    If text_color is None, auto-detects per page (use detect_text_color() for consistency).
    """
    result = _get_watermark_region(image_path)
    if result is None:
        Image.open(image_path).save(output_path)
        return False

    image, r1, r2, c1, c2 = result
    region_w = c2 - c1
    region_h = r2 - r1

    if image.mode != "RGBA":
        image = image.convert("RGBA")

    image_array = np.array(image)
    bg_color, _ = _detect_bg(image_array, r1, r2, c1, c2)

    # Fill watermark region with background color
    image_array[r1:r2, c1:c2] = bg_color
    image = Image.fromarray(image_array)

    if not text:
        # No text — just the clean erase
        image.convert("RGB").save(output_path)
        return True

    # Render text stamp
    font_size = max(16, int(region_h * 0.55))
    font = _load_font(font_size)

    # Use provided color or auto-detect
    if text_color is None:
        _, brightness = _detect_bg(image_array, r1, r2, c1, c2)
        if brightness > 127:
            text_color = (80, 80, 80, 200)
        else:
            text_color = (200, 200, 200, 200)

    # Create text overlay
    overlay = Image.new("RGBA", image.size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(overlay)

    bbox = draw.textbbox((0, 0), text, font=font)
    tw = bbox[2] - bbox[0]
    th = bbox[3] - bbox[1]

    # Center text in the watermark region
    tx = c1 + (region_w - tw) // 2
    ty = r1 + (region_h - th) // 2 - bbox[1]  # adjust for font ascent

    draw.text((tx, ty), text, fill=text_color, font=font)
    image = Image.alpha_composite(image, overlay)

    image.convert("RGB").save(output_path)
    return True


def remove_watermark(image_path: Path, output_path: Path) -> bool:
    """Remove NotebookLM watermark from a single image using biharmonic inpainting."""
    result = _get_watermark_region(image_path)
    if result is None:
        Image.open(image_path).save(output_path)
        return False

    image, r1, r2, c1, c2 = result
    image_array = np.array(image)

    mask = np.zeros(image_array.shape[:2], dtype=bool)
    mask[r1:r2, c1:c2] = True

    inpainted = inpaint.inpaint_biharmonic(image_array, mask, channel_axis=-1)
    Image.fromarray((inpainted * 255).astype("uint8")).save(output_path)
    return True


def images_to_pdf(image_paths: list[Path], output_path: Path):
    """Combine images back into a PDF."""
    doc = fitz.open()
    for img_path in sorted(image_paths):
        img = fitz.open(str(img_path))
        rect = img[0].rect
        pdf_page = doc.new_page(width=rect.width, height=rect.height)
        pdf_page.insert_image(rect, filename=str(img_path))
        img.close()
    doc.save(str(output_path))
    doc.close()


def images_to_pptx(image_paths: list[Path], output_path: Path):
    """Create a PowerPoint from cleaned images."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for img_path in sorted(image_paths):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        slide.shapes.add_picture(
            str(img_path), Inches(0), Inches(0),
            width=prs.slide_width, height=prs.slide_height,
        )
    prs.save(str(output_path))


def images_to_zip(image_paths: list[Path], output_path: Path):
    """Package cleaned PNGs into a ZIP."""
    with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for img_path in sorted(image_paths):
            zf.write(img_path, img_path.name)


def process_pdf(
    pdf_path: str,
    output: str | None = None,
    fmt: str = "pdf",
    dpi: int = DPI,
    mode: str = "stamp",
    stamp: str | None = None,
) -> Path:
    """Full pipeline: PDF -> extract pages -> remove watermarks -> reassemble.

    Args:
        mode: "stamp" (page numbers/custom text, default) or "inpaint" (biharmonic erase).
        stamp: Custom text for stamp mode. None = auto page numbers ("1 / N").
    """
    pdf_path = Path(pdf_path)
    if not pdf_path.exists():
        raise FileNotFoundError(f"Not found: {pdf_path}")

    # Determine output path
    suffix_map = {"pdf": ".pdf", "png": ".zip", "pptx": ".pptx"}
    if output:
        out_path = Path(output)
    else:
        out_path = pdf_path.with_stem(pdf_path.stem + "_clean").with_suffix(suffix_map[fmt])

    # Extract pages
    print(f"Extracting pages from {pdf_path.name}...")
    raw_images = pdf_to_images(str(pdf_path), dpi=dpi)
    total = len(raw_images)
    print(f"  {total} pages extracted")

    # Remove watermarks
    clean_dir = Path(tempfile.mkdtemp(prefix="clean_"))
    clean_images = []

    # Detect text color once from the first page for consistency across all pages
    consistent_color = detect_text_color(raw_images[0]) if mode == "stamp" else None

    for i, img in enumerate(raw_images):
        clean_path = clean_dir / img.name
        if mode == "stamp":
            if stamp:
                text = stamp
            elif total > 1:
                text = f"{i + 1} / {total}"
            else:
                # Single page (infographic) — stamp current month/year
                now = datetime.date.today()
                text = f"{now.strftime('%B')} / {now.year}"
            removed = stamp_text(img, clean_path, text=text, text_color=consistent_color)
            status = "stamped" if removed else "no watermark region"
        else:
            removed = remove_watermark(img, clean_path)
            status = "watermark removed" if removed else "no watermark region"
        print(f"  Page {i+1}/{total}: {status}")
        clean_images.append(clean_path)

    # Reassemble
    if fmt == "pdf":
        images_to_pdf(clean_images, out_path)
    elif fmt == "pptx":
        images_to_pptx(clean_images, out_path)
    elif fmt == "png":
        images_to_zip(clean_images, out_path)

    print(f"Output: {out_path}")
    return out_path


def main():
    parser = argparse.ArgumentParser(
        description="Remove NotebookLM watermark from PDF slides/infographics"
    )
    parser.add_argument("input", nargs="+", help="PDF file(s) to process")
    parser.add_argument("-o", "--output", help="Output file path (single file mode)")
    parser.add_argument("--png", action="store_true", help="Output as PNG ZIP")
    parser.add_argument("--pptx", action="store_true", help="Output as PowerPoint")
    parser.add_argument("--dpi", type=int, default=DPI, help=f"Render DPI (default {DPI})")
    parser.add_argument("--inpaint", action="store_true",
                        help="Use biharmonic inpainting instead of stamp")
    parser.add_argument("--stamp", type=str, default=None, metavar="TEXT",
                        help='Custom stamp text (default: page numbers "1 / N")')
    args = parser.parse_args()

    fmt = "pptx" if args.pptx else "png" if args.png else "pdf"
    mode = "inpaint" if args.inpaint else "stamp"

    for pdf_file in args.input:
        try:
            out = args.output if len(args.input) == 1 else None
            process_pdf(pdf_file, output=out, fmt=fmt, dpi=args.dpi,
                        mode=mode, stamp=args.stamp)
        except Exception as e:
            print(f"ERROR processing {pdf_file}: {e}", file=sys.stderr)


if __name__ == "__main__":
    main()
